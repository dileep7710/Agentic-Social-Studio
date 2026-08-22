import os
import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# 16:9 Widescreen dimensions
SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

# Color Palette (Dark Futuristic Theme)
BG_COLOR = RGBColor(15, 16, 38)       # Deep Navy #0F1026
CARD_BG = RGBColor(26, 32, 60)        # Glass Navy #1A203C
ACCENT_CYAN = RGBColor(56, 189, 248)  # Cyan #38BDF8
ACCENT_PURPLE = RGBColor(168, 85, 247)# Purple #A855F7
ACCENT_GOLD = RGBColor(251, 191, 36)  # Amber Gold #FBBF24
TEXT_WHITE = RGBColor(248, 250, 252)  # Bright White
TEXT_MUTED = RGBColor(148, 163, 184)  # Muted Slate #94A3B8
BORDER_COLOR = RGBColor(99, 102, 241) # Indigo #6366F1

prs = Presentation()
prs.slide_width = SLIDE_WIDTH
prs.slide_height = SLIDE_HEIGHT
blank_layout = prs.slide_layouts[6]

def apply_background(slide):
    # Add dark background rectangle
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_WIDTH, SLIDE_HEIGHT)
    bg.fill.solid()
    bg.fill.fore_color.rgb = BG_COLOR
    bg.line.fill.background()
    return bg

def add_header(slide, title_text, category_text="AGENTIC AI OMNI-STUDIO"):
    # Category tag
    cat_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.7), Inches(0.35))
    tf_cat = cat_box.text_frame
    tf_cat.word_wrap = True
    p_cat = tf_cat.paragraphs[0]
    p_cat.text = f"● {category_text.upper()}"
    p_cat.font.size = Pt(11)
    p_cat.font.bold = True
    p_cat.font.color.rgb = ACCENT_CYAN

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.75), Inches(11.7), Inches(0.8))
    tf_title = title_box.text_frame
    tf_title.word_wrap = True
    p_title = tf_title.paragraphs[0]
    p_title.text = title_text
    p_title.font.size = Pt(24)
    p_title.font.bold = True
    p_title.font.color.rgb = TEXT_WHITE

def add_card(slide, left, top, width, height, title, points, accent=ACCENT_CYAN):
    # Card shape
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    card.fill.solid()
    card.fill.fore_color.rgb = CARD_BG
    card.line.color.rgb = BORDER_COLOR
    card.line.width = Pt(1.5)

    # Content
    tb = slide.shapes.add_textbox(left + Inches(0.25), top + Inches(0.2), width - Inches(0.5), height - Inches(0.4))
    tf = tb.text_frame
    tf.word_wrap = True

    if title:
        p0 = tf.paragraphs[0]
        p0.text = title
        p0.font.size = Pt(16)
        p0.font.bold = True
        p0.font.color.rgb = accent
        p0.space_after = Pt(10)

    for pt in points:
        p = tf.add_paragraph()
        p.text = f"• {pt}"
        p.font.size = Pt(13)
        p.font.color.rgb = TEXT_WHITE
        p.space_after = Pt(6)

def add_notes(slide, notes_text):
    notes_slide = slide.notes_slide
    text_frame = notes_slide.notes_text_frame
    text_frame.text = f"SPEAKER NOTES (बोलने के लिए):\n\n{notes_text}"


# ==========================================
# SLIDE 1: TITLE SLIDE
# ==========================================
s1 = prs.slides.add_slide(blank_layout)
apply_background(s1)

tbox = s1.shapes.add_textbox(Inches(1.2), Inches(1.8), Inches(10.9), Inches(3.8))
tf = tbox.text_frame
tf.word_wrap = True

p1 = tf.paragraphs[0]
p1.text = "🌌 AGENTIC AI OMNI-STUDIO"
p1.font.size = Pt(36)
p1.font.bold = True
p1.font.color.rgb = ACCENT_GOLD
p1.space_after = Pt(10)

p2 = tf.add_paragraph()
p2.text = "Autonomous Multi-Platform Content Orchestration with AES-256-GCM Security & 4K Visual Engine"
p2.font.size = Pt(18)
p2.font.color.rgb = ACCENT_CYAN
p2.space_after = Pt(25)

p3 = tf.add_paragraph()
p3.text = "Author: Dileep Yadav  |  Department: Computer Science & Engineering (B.Tech)  |  2026"
p3.font.size = Pt(14)
p3.font.color.rgb = TEXT_MUTED

add_notes(s1, "Respected professors and jury members, good morning. Aaj main apna B.Tech capstone project present kar raha hoon jiska naam hai 'Agentic AI Omni-Studio'. Yeh ek autonomous AI content orchestration aur secure broadcasting system hai, jo 1 single creative thought ko 5 alag-alag social media platforms ke hisaab se customize karta hai aur 1-click me publish karta hai.")

# ==========================================
# SLIDE 2: PROBLEM STATEMENT
# ==========================================
s2 = prs.slides.add_slide(blank_layout)
apply_background(s2)
add_header(s2, "Problem Statement & Industry Challenges", "CHALLENGES IN SOCIAL AUTOMATION")
add_card(s2, Inches(0.8), Inches(1.8), Inches(5.6), Inches(5.0), "1. Content Fragmentation", [
    "Every platform requires unique copy, tone, and formatting.",
    "LinkedIn demands professional executive structure.",
    "Instagram requires aesthetic emojis and visual hashtags.",
    "X (Twitter) enforces strict 280-character viral limits.",
    "Copy-pasting identical text causes audience disconnect."
], ACCENT_GOLD)
add_card(s2, Inches(6.8), Inches(1.8), Inches(5.6), Inches(5.0), "2. Friction & Security Risks", [
    "Hootsuite / Buffer require complex developer accounts and fees.",
    "Normal users cannot configure OAuth tokens or Facebook Apps.",
    "Traditional automation tools store access tokens in plaintext.",
    "Exposed tokens lead to account hijacking and data leaks.",
    "Creators waste 2-3 hours daily doing repetitive resizing."
], ACCENT_PURPLE)
add_notes(s2, "Social media management me sabse badi problem fragmentation ki hai. Har platform ki audience aur format alag hai. Plus traditional tools bohot complex hain aur tokens plain text me store karke security risk banate hain.")

# ==========================================
# SLIDE 3: PROJECT VISION & OBJECTIVES
# ==========================================
s3 = prs.slides.add_slide(blank_layout)
apply_background(s3)
add_header(s3, "Project Vision & Core Objectives", "CORE GOALS")
add_card(s3, Inches(0.8), Inches(1.8), Inches(3.6), Inches(5.0), "Autonomous Planning", [
    "Goal extraction from raw ideas.",
    "Sequential sub-task creation.",
    "Self-correcting error handling.",
    "Works offline & online seamlessly."
], ACCENT_CYAN)
add_card(s3, Inches(4.8), Inches(1.8), Inches(3.6), Inches(5.0), "Content & 4K Visuals", [
    "1 Input -> 5 Distinct Copies (Slide 8).",
    "4K Frosted Glassmorphism rendering.",
    "TrueType scalable typography.",
    "Custom signature watermark."
], ACCENT_GOLD)
add_card(s3, Inches(8.8), Inches(1.8), Inches(3.6), Inches(5.0), "Enterprise Security", [
    "AES-256-GCM token encryption at rest.",
    "HttpOnly short-lived session cookies.",
    "Refresh Token Rotation (Remember Me).",
    "BOLA / IDOR strict authorization."
], ACCENT_PURPLE)
add_notes(s3, "Hamara objective ek aisa system banana tha jo AI autonomy, enterprise-grade security aur consumer-level simplicity ko ek sath laaye — taaki chahe technical user ho ya non-technical, sab 1-click me broadcast kar sakein.")

# ==========================================
# SLIDE 4: SYSTEM ARCHITECTURE
# ==========================================
s4 = prs.slides.add_slide(blank_layout)
apply_background(s4)
add_header(s4, "5-Tier High-Level System Architecture", "SYSTEM ARCHITECTURE")
add_card(s4, Inches(0.8), Inches(1.8), Inches(11.7), Inches(5.0), "End-to-End Orchestration Pipeline", [
    "Tier 1 (Client Layer): Streamlit Interactive Omni-Studio + React 18 / Tailwind SaaS Dashboard.",
    "Tier 2 (Gateway & Security): FastAPI with CSP headers, sliding-window rate limiters, and HttpOnly cookies.",
    "Tier 3 (Cognitive AI Core): Goal Manager, AI Planner (Llama 3.2 / Gemini), and Platform Content Adapter.",
    "Tier 4 (Visual Engine & CDN): Pillow 4K Frosted Glass Renderer + Catbox / TmpFiles Multi-CDN uploader.",
    "Tier 5 (Data & Security Tier): SQLite with UUID session isolation + AES-256-GCM encrypted credentials."
], ACCENT_CYAN)
add_notes(s4, "Hamara system 5-tier architecture par based hai. Client request FastAPI gateway par aati hai jo rate-limiter aur security headers se protected hai. Wahan se AI Core content plan karta hai, Pillow engine 4K media generate karta hai, aur dual-mode dispatchers use broadcast karte hain.")

# ==========================================
# SLIDE 5: AGENTIC AI CORE (8 MODULES)
# ==========================================
s5 = prs.slides.add_slide(blank_layout)
apply_background(s5)
add_header(s5, "Agentic AI Cognitive Core (8 Modules)", "PPT SLIDES 7, 9, 14 ALIGNMENT")
add_card(s5, Inches(0.8), Inches(1.8), Inches(5.6), Inches(5.0), "Perception & Planning Modules", [
    "Module 1 (Goal Manager): Theme classification & intent extraction.",
    "Module 2 (AI Planner): Task decomposition into 4 steps.",
    "Module 3 (Content Adapter): Platform-specific copy generation.",
    "Module 4 (Tool Selector): Automatic credentials & protocol selection."
], ACCENT_GOLD)
add_card(s5, Inches(6.8), Inches(1.8), Inches(5.6), Inches(5.0), "Execution & Recovery Modules", [
    "Module 5 (Execution Engine): Parallel multi-destination dispatchers.",
    "Module 6 (Evaluation Engine): Partial-success resilience tracking.",
    "Module 7 (Memory System): Session-isolated SQLite persistence.",
    "Module 8 (Autonomous Fallback): API failure to 1-Tap Intent failover."
], ACCENT_PURPLE)
add_notes(s5, "Simple chatbots aur Agentic AI me yeh farak hai ki Agentic AI ek cognitive loop follow karta hai — yeh goal samajhta hai, plan banata hai, appropriate tool choose karta hai aur error aane par khud fallback strategy apply karta hai.")

# ==========================================
# SLIDE 6: GOAL UNDERSTANDING & PLANNING
# ==========================================
s6 = prs.slides.add_slide(blank_layout)
apply_background(s6)
add_header(s6, "Goal Understanding & Sequential Planning", "MODULE 1 & 2 DEEP DIVE")
add_card(s6, Inches(0.8), Inches(1.8), Inches(11.7), Inches(5.0), "Deterministic + Neural Hybrid Planning", [
    "1. Intent Analysis: Detects themes like 'Success', 'Innovation', 'Mindset' from user prompt.",
    "2. Step 1 (Theme & Context): Extracts key takeaways and hashtags suitable for the tone.",
    "3. Step 2 (Visual Asset Generation): Renders 4K aesthetic background canvas with signature.",
    "4. Step 3 (Copy Adaptation): Formulates 5 distinct copywriting variations tailored to each network.",
    "5. Step 4 (Multi-Channel Dispatch): Broadcasts to active channels with individual confirmation.",
    "6. Offline Resilience: If local LLM (Ollama) is unavailable, deterministic expert heuristics take over instantly."
], ACCENT_CYAN)
add_notes(s6, "Goal Manager user ke input ka theme pehchanta hai aur AI Planner use 4 sequential sub-tasks me divide karta hai. Isme humne deterministic fallback bhi diya hai taaki internet na hone par bhi planning 100% chale.")

# ==========================================
# SLIDE 7: CONTENT ADAPTATION ENGINE (SLIDE 8)
# ==========================================
s7 = prs.slides.add_slide(blank_layout)
apply_background(s7)
add_header(s7, "Multi-Platform Content Adaptation (Slide 8)", "1 INPUT -> 5 PLATFORM COPIES")
add_card(s7, Inches(0.8), Inches(1.8), Inches(5.6), Inches(5.0), "Visual & Mobile Networks", [
    "📸 Instagram: High aesthetic visual appeal + 15 curated lifestyle hashtags (#MindsetMatters, #DailyInspiration).",
    "💬 WhatsApp: Clean markdown (*Bold*, _Italics_) + inspiring message + direct 4K media link.",
    "🐦 X (Twitter): Crisp viral tweet strictly capped under 280 characters with relevant hashtags."
], ACCENT_GOLD)
add_card(s7, Inches(6.8), Inches(1.8), Inches(5.6), Inches(5.0), "Professional & Social Networks", [
    "💼 LinkedIn: Formal executive structure, key takeaway bullet points, leadership hashtags (#Leadership, #Growth).",
    "📘 Facebook: Community engagement questions ('What are your thoughts on this?') with discussion prompts."
], ACCENT_PURPLE)
add_notes(s7, "Slide 8 concept hamara core USP hai. Content Adapter ek hi quote ko 5 alag formats me convert karta hai taaki har social media algorithm par post organically perform kare bina copy-paste lage.")

# ==========================================
# SLIDE 8: 4K VISUAL GENERATION ENGINE
# ==========================================
s8 = prs.slides.add_slide(blank_layout)
apply_background(s8)
add_header(s8, "4K Visual Generation & Glassmorphism", "PILLOW RENDERING ENGINE")
add_card(s8, Inches(0.8), Inches(1.8), Inches(11.7), Inches(5.0), "Programmatic Aesthetic Canvas Construction", [
    "Dimensions: 1080x1920 (9:16 Story format) & 1080x1080 (1:1 Feed format).",
    "Layer 1 (Background): High-resolution curated nature wallpapers (Galaxy, Mountain, Sunset).",
    "Layer 2 (Frosted Glass Card): Semi-transparent RGBA dark card (15, 23, 42, 205) with 36px rounded corners.",
    "Layer 3 (Typography): Scalable TrueType font rendering with dynamic line wrapping and anchor centering.",
    "Layer 4 (Watermark Signature): Golden accent (-- Signature Name) embedded cleanly at the bottom.",
    "Performance: Complete 4K render in under 0.78 seconds with 95% JPEG/PNG compression quality."
], ACCENT_CYAN)
add_notes(s8, "Hamara visual engine Pillow library use karke 0.78 seconds me multi-layer frosted glass card render karta hai. Yeh dynamically text width aur line spacing calculate karke user ka custom watermark embed karta hai.")

# ==========================================
# SLIDE 9: MULTI-CDN ASSET PIPELINE
# ==========================================
s9 = prs.slides.add_slide(blank_layout)
apply_background(s9)
add_header(s9, "High-Speed Multi-CDN Asset Pipeline", "PUBLIC MEDIA URL GENERATION")
add_card(s9, Inches(0.8), Inches(1.8), Inches(11.7), Inches(5.0), "Automated High-Availability Upload", [
    "Problem: Meta Graph API & LinkedIn API reject local C:\\ file paths; they require publicly accessible HTTPS URLs.",
    "Primary Server: Catbox.moe High-Speed File API with automated file multipart upload.",
    "Secondary Server: TmpFiles Backup CDN activated automatically if Primary fails.",
    "Timeout Protection: 30-second timeout guard with automatic retry mechanism.",
    "Result: Generates permanent, public HTTPS URLs (e.g. https://files.catbox.moe/abc123.png) in ~1.2 seconds."
], ACCENT_GOLD)
add_notes(s9, "Meta aur LinkedIn APIs ko image download karne ke liye public HTTPS URL chahiye hota hai. Hamara multi-CDN pipeline local image ko automatically upload karke instant live link provide karta hai.")

# ==========================================
# SLIDE 10: DUAL-MODE DISPATCH ARCHITECTURE
# ==========================================
s10 = prs.slides.add_slide(blank_layout)
apply_background(s10)
add_header(s10, "Dual-Mode Dispatch Architecture", "ZERO-TOKEN VS OAUTH API")
add_card(s10, Inches(0.8), Inches(1.8), Inches(5.6), Inches(5.0), "Mode 1: Normal Users (Zero-Token)", [
    "1-Tap Native Mobile Share Sheet (navigator.share).",
    "WhatsApp Click-to-Chat URI Intent (wa.me).",
    "Twitter/X 1-Click Tweet Intent.",
    "Facebook Timeline Sharer Intent.",
    "0% developer tokens, 0% passwords, 0% ban risk."
], ACCENT_CYAN)
add_card(s10, Inches(6.8), Inches(1.8), Inches(5.6), Inches(5.0), "Mode 2: Pro Creators (OAuth API)", [
    "Meta Graph API v21 Direct Container Publishing.",
    "LinkedIn REST API v2 UGC Post Publishing.",
    "Automatic background broadcast without opening apps.",
    "Partial-Success Engine (1 failure doesn't block others)."
], ACCENT_PURPLE)
add_notes(s10, "Humne dual-mode dispatch banaya hai. Normal users Canva aur Spotify ki tarah bina password/token ke 1-click me share kar sakte hain. Aur professional creators Meta aur LinkedIn tokens daal kar background me automated post kar sakte hain.")

# ==========================================
# SLIDE 11: AES-256-GCM TOKEN ENCRYPTION
# ==========================================
s11 = prs.slides.add_slide(blank_layout)
apply_background(s11)
add_header(s11, "Enterprise Security — AES-256-GCM Encryption", "TOKEN ENCRYPTION AT REST")
add_card(s11, Inches(0.8), Inches(1.8), Inches(11.7), Inches(5.0), "Authenticated Encryption with Associated Data (AEAD)", [
    "Cipher: AES-256 in Galois/Counter Mode (GCM) via cryptography.hazmat library.",
    "Cryptographic Nonce: 12-byte cryptographically secure random IV generated per token.",
    "Storage Format: aes_gcm:v1:<nonce_hex>:<ciphertext_and_tag_hex>.",
    "Integrity Verification: 16-byte authentication tag prevents bit-flipping and padding oracle attacks.",
    "Key Management: 256-bit master key derived via SHA-256 from secure environment variable.",
    "Zero Plaintext Guarantee: Database contains 0 unencrypted OAuth tokens."
], ACCENT_GOLD)
add_notes(s11, "Security ke liye humne sabhi OAuth tokens ko AES-256-GCM se encrypt kiya hai. Agar koi hacker database file chori bhi kar le, toh bhi wo raw tokens ko decrypt nahi kar payega.")

# ==========================================
# SLIDE 12: AUTHENTICATION & REFRESH ROTATION
# ==========================================
s12 = prs.slides.add_slide(blank_layout)
apply_background(s12)
add_header(s12, "Authentication & Refresh Token Rotation", "SESSION MANAGEMENT")
add_card(s12, Inches(0.8), Inches(1.8), Inches(5.6), Inches(5.0), "Password Hashing & Cookies", [
    "PBKDF2-HMAC-SHA256 with 150,000 iterations.",
    "16-byte random cryptographic salt.",
    "Zero plaintext fallback comparison.",
    "Delivered via HttpOnly + Secure + SameSite=Lax cookies (100% XSS protected)."
], ACCENT_CYAN)
add_card(s12, Inches(6.8), Inches(1.8), Inches(5.6), Inches(5.0), "Refresh Token Rotation", [
    "Access Token: 15 minutes short-lived validity.",
    "Refresh Token: 30 days (Remember Me persistent).",
    "On every refresh, old refresh token is revoked.",
    "Issuing a fresh pair eliminates Replay Attacks."
], ACCENT_PURPLE)
add_notes(s12, "Humne localStorage ke bajaye HttpOnly cookies use ki hain jo XSS attacks se bachati hain. Refresh Token Rotation se user Google/Instagram ki tarah securely logged in rehta hai bina bar-bar password dale.")

# ==========================================
# SLIDE 13: BRUTE-FORCE DEFENSE & HEADERS
# ==========================================
s13 = prs.slides.add_slide(blank_layout)
apply_background(s13)
add_header(s13, "Brute-Force Defense & Security Headers", "THREAT MITIGATION")
add_card(s13, Inches(0.8), Inches(1.8), Inches(11.7), Inches(5.0), "Multi-Layer Attack Prevention", [
    "Sliding-Window IP Rate Limiter: Max 10 attempts per 15 minutes on authentication endpoints.",
    "Account Lockout: 5 consecutive failed login attempts trigger HTTP 423 Locked for 15 minutes.",
    "Content-Security-Policy (CSP): Restricts script execution to self and trusted domains.",
    "X-Frame-Options: DENY (Prevents Clickjacking attacks).",
    "X-Content-Type-Options: nosniff (Prevents MIME-type sniffing).",
    "Referrer-Policy: strict-origin-when-cross-origin."
], ACCENT_GOLD)
add_notes(s13, "Credential stuffing aur brute-force attacks ko rokne ke liye humne 5-attempt lockout policy lagayi hai, sath hi Content Security Policy aur strict HTTP headers implement kiye hain.")

# ==========================================
# SLIDE 14: MEDIA MAGIC-BYTE VALIDATION
# ==========================================
s14 = prs.slides.add_slide(blank_layout)
apply_background(s14)
add_header(s14, "Media Anti-Tamper & Magic-Byte Validation", "FILE UPLOAD SECURITY")
add_card(s14, Inches(0.8), Inches(1.8), Inches(11.7), Inches(5.0), "Header Byte Inspection vs Extension Trust", [
    "Vulnerability Mitigated: Renaming malicious shell scripts or executables to fake .png extensions.",
    "Magic Byte Verification: Inspects first 64 raw bytes of uploaded files:",
    "   • PNG Signature: \\x89PNG\\r\\n\\x1a\\n",
    "   • JPEG Signature: \\xff\\xd8\\xff",
    "   • WebP Signature: RIFF....WEBP",
    "   • MP4 Signature: ftyp box header in first 20 bytes",
    "Payload Rejection: Non-matching files rejected with HTTP 400 Bad Request.",
    "Size Quota: Strict 25 MB file size limit enforced."
], ACCENT_CYAN)
add_notes(s14, "Hum file extensions par trust nahi karte. Server file ke starting 64 bytes (Magic Bytes) inspect karta hai taaki koi malicious script ya executable ko .png bana kar upload na kar sake.")

# ==========================================
# SLIDE 15: MULTI-USER ISOLATION & IDOR DEFENSE
# ==========================================
s15 = prs.slides.add_slide(blank_layout)
apply_background(s15)
add_header(s15, "Multi-User Isolation & IDOR / BOLA Defense", "AUTHORIZATION BARRIERS")
add_card(s15, Inches(0.8), Inches(1.8), Inches(11.7), Inches(5.0), "Strict Server-Side Ownership Enforcement", [
    "Broken Object Level Authorization (BOLA / IDOR) Defense:",
    "   • Client-supplied user_id parameters are strictly ignored and never trusted.",
    "   • Database queries filter by verified server session: WHERE user_id = current_user.id.",
    "   • Unauthorized cross-account operations return HTTP 403 Forbidden.",
    "Streamlit Session Isolation:",
    "   • Every visitor receives a UUID-isolated session_id in SQLite.",
    "   • Guarantees 0% data bleed between concurrent users on public deployments."
], ACCENT_PURPLE)
add_notes(s15, "IDOR vulnerability se bachne ke liye client ke bheje hue user_id par kabhi trust nahi kiya jata; server cryptographically signed token se user verify karke hi data access allow karta hai.")

# ==========================================
# SLIDE 16: DATABASE SCHEMA & ENTITIES
# ==========================================
s16 = prs.slides.add_slide(blank_layout)
apply_background(s16)
add_header(s16, "Database Schema & Entity Relationships", "DATA PERSISTENCE MODEL")
add_card(s16, Inches(0.8), Inches(1.8), Inches(11.7), Inches(5.0), "Normalized Relational Architecture (SQLAlchemy)", [
    "1. users: id, name, email, hashed_password, failed_attempts, locked_until, last_login.",
    "2. user_sessions: id, user_id, session_token_hash, refresh_token_hash, device_info, is_revoked.",
    "3. social_accounts: id, user_id, platform, account_name, access_token (AES encrypted), status.",
    "4. post_logs & platform_post_results: Historical audit trail of all omni-channel publications.",
    "5. audit_logs: Structured security events (login, logout, token_refresh) with zero secret leakage.",
    "6. media_files: Validated MIME metadata, file sizes, and storage paths.",
    "GDPR Compliance: DELETE /api/account/me cascades and permanently purges all user data."
], ACCENT_GOLD)
add_notes(s16, "Database schema fully normalized hai aur GDPR 'Right to be Forgotten' rule follow karta hai — account delete karne par user ke sabhi encrypted tokens, media aur sessions instantly wipe ho jate hain.")

# ==========================================
# SLIDE 17: BILINGUAL USER EXPERIENCE
# ==========================================
s17 = prs.slides.add_slide(blank_layout)
apply_background(s17)
add_header(s17, "Bilingual Glassmorphism User Interface", "USER EXPERIENCE & ACCESSIBILITY")
add_card(s17, Inches(0.8), Inches(1.8), Inches(5.6), Inches(5.0), "Streamlit Studio (app.py)", [
    "Real-time 4K live canvas preview.",
    "Instant English / हिन्दी language toggle.",
    "1-Tap Native Mobile Share Sheet.",
    "Smart input validation (2-letter name & 10-15 digit phone auto-formatting)."
], ACCENT_CYAN)
add_card(s17, Inches(6.8), Inches(1.8), Inches(5.6), Inches(5.0), "FullStack SaaS UI (React 18)", [
    "Active Devices Manager (Browser/OS/IP).",
    "1-Click 'Log Out of All Devices'.",
    "Individual device session revocation.",
    "Security audit logs viewer."
], ACCENT_PURPLE)
add_notes(s17, "UI modern Glassmorphism theme par design kiya gaya hai. Isme English aur Hindi dono ka support hai, real-time input verification hai aur active sessions manage karne ka pura control hai.")

# ==========================================
# SLIDE 18: BENCHMARKS & TEST RESULTS
# ==========================================
s18 = prs.slides.add_slide(blank_layout)
apply_background(s18)
add_header(s18, "Testing & Verification Benchmark Results", "23/23 TESTS PASSED (100%)")
add_card(s18, Inches(0.8), Inches(1.8), Inches(5.6), Inches(5.0), "Automated Test Suite Pass Rate", [
    "test_security_suite.py: 8/8 PASSED (PBKDF2, AES-256-GCM, Rotation, Lockout, IDOR, Magic Bytes).",
    "test_production_suite.py: 8/8 PASSED (10 Concurrent users, 4K render, Multi-Account).",
    "test_oauth_suite.py: 7/7 PASSED (Meta OAuth, Signed JWT State, Multi-Account)."
], ACCENT_GOLD)
add_card(s18, Inches(6.8), Inches(1.8), Inches(5.6), Inches(5.0), "Performance Benchmarks", [
    "4K Graphic Render Latency: ~0.78 seconds.",
    "Multi-CDN Upload Latency: ~1.2 seconds.",
    "End-to-End Omni-Broadcast Latency: < 2.5 seconds.",
    "Concurrency: 0 race conditions across 10 parallel threads."
], ACCENT_CYAN)
add_notes(s18, "Humne 23 automated tests run kiye jisme token encryption, brute-force lockout, IDOR defense aur 10 simultaneous users ka load test shamil tha. Saare tests 100% pass hue.")

# ==========================================
# SLIDE 19: FUTURE SCOPE & ROADMAP
# ==========================================
s19 = prs.slides.add_slide(blank_layout)
apply_background(s19)
add_header(s19, "Future Scope & Technical Roadmap", "FUTURE ENHANCEMENTS")
add_card(s19, Inches(0.8), Inches(1.8), Inches(11.7), Inches(5.0), "Upcoming Capabilities", [
    "1. Video Reels Engine: Automated generation of 15-second animated reels with background music.",
    "2. Neural AI Voiceovers: Integration with ElevenLabs / Edge-TTS for multi-lingual spoken quotes.",
    "3. Smart Calendar Scheduler: Cron-based automated posting queue with visual content calendar.",
    "4. Centralized Analytics Dashboard: Aggregated engagement metrics (views, likes, comments, shares).",
    "5. AI Thumbnail A/B Testing: Automated variation generation to maximize social engagement."
], ACCENT_GOLD)
add_notes(s19, "Future me hum isme automated AI Video Reels generator, multi-lingual neural voice synthesis aur calendar-based automated scheduling system add karne ka plan kar rahe hain.")

# ==========================================
# SLIDE 20: CONCLUSION & SUMMARY
# ==========================================
s20 = prs.slides.add_slide(blank_layout)
apply_background(s20)
tbox = s20.shapes.add_textbox(Inches(1.2), Inches(1.8), Inches(10.9), Inches(4.5))
tf = tbox.text_frame
tf.word_wrap = True

p1 = tf.paragraphs[0]
p1.text = "🏆 CONCLUSION & SUMMARY"
p1.font.size = Pt(32)
p1.font.bold = True
p1.font.color.rgb = ACCENT_GOLD
p1.space_after = Pt(15)

points = [
    "Built a production-grade Agentic AI Omni-Studio with 1-click multi-platform broadcasting.",
    "Eliminated copy-paste fragmentation via 5-platform AI content adaptation (Slide 8).",
    "Provided zero-token simplicity for normal users alongside professional OAuth API workflows.",
    "Hardened with AES-256-GCM token encryption, HttpOnly cookies, and Refresh Token Rotation.",
    "Live Deployment: https://dileep-ai-studio.streamlit.app  |  GitHub: Agentic-Social-Studio"
]

for pt in points:
    p = tf.add_paragraph()
    p.text = f"• {pt}"
    p.font.size = Pt(15)
    p.font.color.rgb = TEXT_WHITE
    p.space_after = Pt(8)

p_end = tf.add_paragraph()
p_end.text = "\nThank you! I am now ready for your questions. (Q&A Session)"
p_end.font.size = Pt(16)
p_end.font.bold = True
p_end.font.color.rgb = ACCENT_CYAN

add_notes(s20, "In conclusion, Agentic AI Omni-Studio yeh prove karta hai ki complex AI architecture ko bank-grade security aur extreme simplicity ke sath deploy kiya ja sakta hai. Thank you professors, ab aapke questions ke liye main ready hoon.")

# Save presentation to root directory
output_file = "Agentic_AI_Omni_Studio_Presentation.pptx"
prs.save(output_file)
print(f"PowerPoint Presentation successfully created: {output_file}")
